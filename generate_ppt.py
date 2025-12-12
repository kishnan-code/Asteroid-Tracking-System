from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper to clean up slide creation
    def add_slide(title_text, content_points, layout_index=1):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = title_text
        
        # Content
        if content_points:
            # Check if using the standard bullet layout
            if layout_index == 1:
                tf = slide.placeholders[1].text_frame
                for i, point in enumerate(content_points):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = point
                    p.font.size = Pt(20)
            else:
                # Custom placement for other layouts if needed
                pass
        return slide

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Astronaut Health AI"
    subtitle = slide.placeholders[1]
    subtitle.text = "Predictive Health Monitoring System for Space Missions\nBy Null and Beyond"

    # 2. Problem Statement
    add_slide("Problem Statement", [
        "Space missions expose astronauts to unique health risks:",
        "• High Radiation levels",
        "• Microgravity effects (Muscle atrophy, Bone density loss)",
        "• Psychological stress and Isolation",
        "• Limited onboard medical resources",
        "Need: A proactive system to monitor and predict health status in real-time."
    ])

    # 3. Our Solution
    add_slide("Project Overview", [
        "A Web-based Machine Learning application that:",
        "1. Collects 10 key biometric & environmental parameters",
        "2. Analyzes data using a trained Predictive Model",
        "3. Calculates a comprehensive 'Health Score' (0-100)",
        "4. Provides actionable AI-driven recommendations",
        "5. Visualizes trends via an interactive Dashboard"
    ])

    # 4. Key Features
    add_slide("Key Features", [
        "Real-time Health Scoring: Instant analysis of user inputs.",
        "Comprehensive Metrics: Tracks Respiration, BP, Oxygen, Radiation, etc.",
        "Interactive Dashboard: Visualizes inputs vs. mission averages.",
        "EDA Integration: Exploratory Data Analysis charts for historical context.",
        "Prediction History: Logs previous scans for tracking trends over time.",
        "Cosmic UI: Immersive, space-themed dark interface."
    ])

    # 5. Technology Stack
    add_slide("Technology Stack", [
        "Backend:",
        "• Python (Core Logic)",
        "• Flask (Web Framework)",
        "Machine Learning:",
        "• Scikit-Learn / Joblib (Model Serialization)",
        "• Pandas & NumPy (Data Processing)",
        "Frontend:",
        "• HTML5 & CSS3 (Custom 'Cosmic' Theme)",
        "• Jinja2 Templating",
        "• JavaScript (Dynamic interactions)"
    ])

    # 6. Workflow / Architecture
    add_slide("System Workflow", [
        "1. Input: User enters data via the Mission Form.",
        "2. Process: Flask app receives data & pre-processes it.",
        "3. Predict: ML Model (pkl) computes the Health Score.",
        "4. Analysis: System compares inputs vs. historical averages.",
        "5. Output: Dashboard displays Score, Status, and Charts."
    ])

    # 7. Future Scope
    add_slide("Future Scope", [
        "IoT Integration: Connect directly to wearable biosensors.",
        "Real-time Telemetry: Live data stream from spacecraft systems.",
        "Deep Learning: Use LSTM/RNN for time-series anomaly detection.",
        "Personalized Profiles: Adaptive baselines for individual astronauts."
    ])

    # 8. Conclusion
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Thank You"
    tf = slide.placeholders[1].text_frame
    p = tf.paragraphs[0]
    p.text = "Ready for Questions?"
    p.alignment = PP_ALIGN.CENTER

    prs.save('Astronaut_Health_AI_Presentation.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
